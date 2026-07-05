from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from sidecar.extractors.utils import rows_to_markdown_table, utc_timestamp
from sidecar.models import ContentType, DocumentElement, DocumentMetadata, ElementType, ExtractedDocument


def extract_pptx_document(source_id: str, file_path: str) -> ExtractedDocument:
    path = Path(file_path).expanduser().resolve()
    if not path.exists() or not path.is_file():
        raise FileNotFoundError(f"File not found at {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Unsupported file type: {path.suffix}")

    presentation = Presentation(str(path))
    elements = _extract_pptx_elements(presentation)
    properties = presentation.core_properties
    metadata = DocumentMetadata(
        author=properties.author or None,
        pageCount=len(presentation.slides),
        createdAt=properties.created.isoformat() if properties.created else None,
        extractedAt=utc_timestamp(),
    )
    return ExtractedDocument(
        sourceId=source_id,
        title=properties.title or path.stem,
        language=None,
        contentType=ContentType.pptx,
        metadata=metadata,
        elements=elements,
    )


def _extract_pptx_elements(presentation: Presentation) -> list[DocumentElement]:
    elements: list[DocumentElement] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.has_text_frame and title_shape.text.strip():
            _append_element(
                elements,
                ElementType.heading,
                title_shape.text.strip(),
                slide_number,
                level=2,
            )

        shapes = sorted(slide.shapes, key=lambda shape: (shape.top or 0, shape.left or 0))
        for shape in shapes:
            if shape is title_shape:
                continue
            _extract_shape(elements, shape, slide_number)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip() if notes_frame is not None else ""
            if notes:
                _append_element(
                    elements,
                    ElementType.paragraph,
                    f"[Speaker Notes, Slide {slide_number}]: {notes}",
                    slide_number,
                    metadata={"kind": "speakerNotes"},
                )

    return elements or [DocumentElement(type=ElementType.paragraph, content="", position=0)]


def _extract_shape(elements: list[DocumentElement], shape, slide_number: int) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        nested_shapes = sorted(shape.shapes, key=lambda nested: (nested.top or 0, nested.left or 0))
        for nested_shape in nested_shapes:
            _extract_shape(elements, nested_shape, slide_number)
        return

    if shape.has_table:
        _append_element(elements, ElementType.table, _table_to_markdown(shape.table), slide_number)
        return

    if shape.has_chart:
        _append_element(elements, ElementType.table, _chart_to_markdown(shape.chart), slide_number)
        return

    if _is_picture(shape):
        alt_text = _picture_alt_text(shape)
        _append_element(
            elements,
            ElementType.image,
            f"![{alt_text}](embedded-image)",
            slide_number,
            metadata={"shapeName": shape.name},
        )
        return

    if not shape.has_text_frame:
        return

    paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
    if not paragraphs:
        return
    if _is_list_shape(shape, paragraphs):
        content = "\n".join(
            f"{'  ' * paragraph.level}- {paragraph.text.strip()}" for paragraph in paragraphs
        )
        _append_element(elements, ElementType.list, content, slide_number)
        return
    _append_element(
        elements,
        ElementType.paragraph,
        "\n".join(paragraph.text.strip() for paragraph in paragraphs),
        slide_number,
    )


def _append_element(
    elements: list[DocumentElement],
    element_type: ElementType,
    content: str,
    page: int,
    *,
    level: int | None = None,
    metadata: dict | None = None,
) -> None:
    if not content.strip():
        return
    elements.append(
        DocumentElement(
            type=element_type,
            content=content.strip(),
            page=page,
            level=level,
            metadata=metadata,
            position=len(elements),
        )
    )


def _is_picture(shape) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    return shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")


def _picture_alt_text(shape) -> str:
    try:
        description = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
    except AttributeError:
        description = ""
    return description.strip() or shape.name or "Image"


def _is_list_shape(shape, paragraphs: Iterable) -> bool:
    paragraphs = list(paragraphs)
    if any(paragraph.level > 0 or _has_explicit_bullet(paragraph) for paragraph in paragraphs):
        return True
    if not shape.is_placeholder or len(paragraphs) < 2:
        return False
    return shape.placeholder_format.type in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}


def _has_explicit_bullet(paragraph) -> bool:
    properties = paragraph._p.pPr
    if properties is None:
        return False
    return any(child.tag.rsplit("}", maxsplit=1)[-1] in {"buChar", "buAutoNum"} for child in properties)


def _table_to_markdown(table) -> str:
    return rows_to_markdown_table([cell.text for cell in row.cells] for row in table.rows)


def _chart_to_markdown(chart) -> str:
    title = chart.chart_title.text_frame.text.strip() if chart.has_title else "Chart"
    try:
        categories = [category.label for category in chart.plots[0].categories]
        series = list(chart.series)
        rows = ["| Category | " + " | ".join(item.name for item in series) + " |"]
        rows.append("| --- | " + " | ".join("---" for _ in series) + " |")
        for index, category in enumerate(categories):
            values = [str(item.values[index]) for item in series]
            rows.append("| " + " | ".join([str(category), *values]) + " |")
        return f"### {title}\n\n" + "\n".join(rows)
    except (IndexError, TypeError, ValueError):
        return f"### {title}\n\n[Unsupported chart]"
