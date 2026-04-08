from datetime import datetime, UTC
from pathlib import Path

from docx import Document
from pptx import Presentation

from sidecar.models import (
    ContentType, ExtractedDocument, DocumentElement,
    ElementType, DocumentMetadata
)

SUPPORTED_EXTENSIONS = {
    ".docx": ContentType.docx,
    ".pptx": ContentType.pptx,
}


def extract_office_document(source_id: str, file_path: str) -> ExtractedDocument:
    path = Path(file_path).expanduser().resolve()

    if not path.exists():
        raise FileNotFoundError(f"File not found at {path}")

    content_type = SUPPORTED_EXTENSIONS.get(path.suffix.lower())
    if content_type is None:
        raise ValueError(f"Unsupported file type: {path.suffix}")

    if content_type == ContentType.docx:
        elements, page_count = _extract_docx_elements(path)
    else:
        elements, page_count = _extract_pptx_elements(path)

    metadata = DocumentMetadata(extractedAt=_timestamp(), pageCount=page_count)

    return ExtractedDocument(
        sourceId=source_id,
        title=path.stem,
        language=None,
        contentType=content_type,
        metadata=metadata,
        elements=elements,
    )


def _extract_docx_elements(path: Path) -> tuple[list[DocumentElement], int]:
    doc = Document(str(path))
    elements: list[DocumentElement] = []
    position = 0

    for block in doc.element.body:
        tag = block.tag.split("}")[-1]

        if tag == "p":
            paragraph = next(
                (p for p in doc.paragraphs if p._element is block), None
            )
            if paragraph is None or not paragraph.text.strip():
                continue

            style_name = paragraph.style.name if paragraph.style else ""

            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.split()[-1])
                except ValueError:
                    level = 1
                elements.append(DocumentElement(
                    type=ElementType("heading"),
                    content=paragraph.text.strip(),
                    level=max(1, min(level, 6)),
                    position=position,
                ))
            else:
                elements.append(DocumentElement(
                    type=ElementType("paragraph"),
                    content=paragraph.text.strip(),
                    position=position,
                ))
            position += 1

        elif tag == "tbl":
            table = next(
                (t for t in doc.tables if t._element is block), None
            )
            if table is None:
                continue

            markdown_table = _table_to_markdown(table)
            elements.append(DocumentElement(
                type=ElementType("table"),
                content=markdown_table,
                position=position,
            ))
            position += 1

    # python-docx has no reliable page count - use section count as approximation
    # or just return 1 if not relevant
    page_count = _estimate_docx_page_count(doc)

    return elements, page_count


def _estimate_docx_page_count(doc) -> int:
    # Word stores page count in the core properties if it was saved by Word
    try:
        return doc.core_properties.revision or 1
    except Exception:
        return 1


def _extract_pptx_elements(path: Path) -> tuple[list[DocumentElement], int]:
    prs = Presentation(str(path))
    elements: list[DocumentElement] = []
    position = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        # Slide heading: first title shape
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            title_text = title_shape.text.strip()
            if title_text:
                elements.append(DocumentElement(
                    type=ElementType("heading"),
                    content=title_text,
                    level=2,
                    page=slide_number,
                    position=position,
                ))
                position += 1

        for shape in slide.shapes:
            if shape is title_shape:
                continue

            if shape.has_table:
                markdown_table = _table_to_markdown(shape.table)
                elements.append(DocumentElement(
                    type=ElementType("table"),
                    content=markdown_table,
                    page=slide_number,
                    position=position,
                ))
                position += 1

            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Detect list-like text frames (bullet points)
                paragraphs = shape.text_frame.paragraphs
                is_list = any(
                    p.text.strip() and p.level > 0
                    for p in paragraphs
                )

                if is_list:
                    bullet_lines = [
                        ("  " * p.level + "- " + p.text.strip())
                        for p in paragraphs
                        if p.text.strip()
                    ]
                    elements.append(DocumentElement(
                        type=ElementType.list,
                        content="\n".join(bullet_lines),
                        page=slide_number,
                        position=position,
                    ))
                else:
                    elements.append(DocumentElement(
                        type=ElementType.paragraph,
                        content=text,
                        page=slide_number,
                        position=position,
                    ))
                position += 1

        # Speaker notes as separate paragraph element
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                elements.append(DocumentElement(
                    type=ElementType.paragraph,
                    content=f"[Speaker Notes, Slide {slide_number}]: {notes_text}",
                    page=slide_number,
                    position=position,
                ))
                position += 1

    return elements, len(prs.slides)


def _table_to_markdown(table) -> str:
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("| " + " | ".join("---" for _ in cells) + " |")
    return "\n".join(rows)


def _timestamp() -> str:
    return datetime.now(UTC).isoformat()
