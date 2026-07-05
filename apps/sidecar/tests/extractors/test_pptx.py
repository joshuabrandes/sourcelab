from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from sidecar.extractors.pptx import extract_pptx_document
from sidecar.models import ContentType, ElementType


def create_presentation(path, image_path):
    presentation = Presentation()
    presentation.core_properties.author = "Grace Hopper"
    presentation.core_properties.title = "Compiler Design"
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "First Slide"
    body = slide.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = "First point"
    body.add_paragraph().text = "Second point"

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "A"
    table_shape.table.cell(0, 1).text = "B"
    table_shape.table.cell(1, 0).text = "1"
    table_shape.table.cell(1, 1).text = "2"
    slide.shapes.add_picture(str(image_path), Inches(6), Inches(1), width=Inches(1))
    slide.notes_slide.notes_text_frame.text = "Explain the implementation."
    presentation.save(path)


def test_pptx_preserves_slide_structure_and_content(tmp_path):
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (20, 20), "white").save(image_path)
    path = tmp_path / "slides.pptx"
    create_presentation(path, image_path)

    extracted = extract_pptx_document("source-1", str(path))

    assert extracted.contentType == ContentType.pptx
    assert extracted.title == "Compiler Design"
    assert extracted.metadata.author == "Grace Hopper"
    assert extracted.metadata.pageCount == 1
    assert all(element.page == 1 for element in extracted.elements)
    assert [element.position for element in extracted.elements] == list(range(len(extracted.elements)))
    assert ElementType.heading in [element.type for element in extracted.elements]
    assert ElementType.list in [element.type for element in extracted.elements]
    assert ElementType.table in [element.type for element in extracted.elements]
    assert ElementType.image in [element.type for element in extracted.elements]
    assert any(element.metadata == {"kind": "speakerNotes"} for element in extracted.elements)
